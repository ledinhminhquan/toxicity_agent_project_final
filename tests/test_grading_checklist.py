from __future__ import annotations

import json
import tempfile
from pathlib import Path

from reportlab.pdfgen import canvas
from pptx import Presentation

from toxicity_agent.autoreport.artifact_loader import ArtifactBundle
from toxicity_agent.grading.checklist import run_grading_checklist


def _make_one_page_pdf(path: Path) -> None:
    c = canvas.Canvas(str(path))
    c.drawString(72, 720, "Test PDF")
    c.showPage()
    c.save()


def _make_one_slide_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test PPTX"
    prs.save(str(path))


def test_grading_checklist_happy_path():
    repo_root = Path(__file__).resolve().parents[1]

    with tempfile.TemporaryDirectory() as d:
        sub = Path(d) / "submission"
        sub.mkdir(parents=True, exist_ok=True)

        report = sub / "report.pdf"
        slides = sub / "slides.pptx"
        manifest = sub / "submission_manifest.json"
        snapshot = sub / "artifacts_snapshot.json"

        _make_one_page_pdf(report)
        _make_one_slide_pptx(slides)
        manifest.write_text(json.dumps({"hello": "world"}), encoding="utf-8")
        snapshot.write_text(json.dumps({"paths": {}}), encoding="utf-8")

        # Create a bundle zip with required files
        import zipfile
        bundle_zip = sub / "submission_bundle.zip"
        with zipfile.ZipFile(bundle_zip, "w", compression=zipfile.ZIP_DEFLATED) as z:
            z.write(report, arcname="report.pdf")
            z.write(slides, arcname="slides.pptx")
            z.write(manifest, arcname="submission_manifest.json")
            z.write(snapshot, arcname="artifacts_snapshot.json")

        bundle = ArtifactBundle(
            train_config={"dataset": {"name": "dummy", "label_fields": ["toxic"]}},
            infer_config={},
            eval_results={
                "baseline_tfidf_lr": {"f1_micro": 0.5},
                "baseline_detoxify_unbiased": {"f1_micro": 0.6},
                "finetuned_transformer": {"f1_micro": 0.7},
            },
            tune_results={"best": {"params": {"learning_rate": 1e-5}}},
            benchmark={"stats": {"p95_ms": 12.3, "p50_ms": 5.0, "p99_ms": 20.0}},
            error_analysis={"top_error_cases": [{"sha256": "abc"}]},
            fairness={"slices": [{"slice_name": "no_identity_mention", "n": 1}]},
            model_metadata={"run_id": "run-123", "model_name": "dummy"},
            thresholds={"best_global": {"threshold": 0.5}, "saved_to": "x"},
            paths={},
        )

        checklist = run_grading_checklist(
            repo_root=repo_root,
            submission_dir=sub,
            bundle=bundle,
            expected_report_pages=(1, 1),
            expected_slides=(1, 1),
        )

        assert checklist["overall_status"] in {"PASS", "WARN"}
        assert checklist["score_estimate_100"] > 0
        ids = {it["id"] for it in checklist["items"]}
        assert "SUB_REPORT" in ids
        assert "SUB_SLIDES" in ids
        assert "REPO_STRUCTURE" in ids
