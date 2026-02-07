from __future__ import annotations

import json
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from ..logging_utils import get_logger
from .artifact_loader import ArtifactBundle
from .charts import (
    save_agent_flow_diagram,
    save_architecture_diagram,
    save_auc_per_label_chart,
    save_eval_comparison_chart,
    save_fairness_gap_chart,
    save_latency_chart,
)

logger = get_logger(__name__)


def _safe(v: Any, default: str = "N/A") -> str:
    if v is None:
        return default
    return str(v)


def _fmt_float(v: Any, digits: int = 3) -> str:
    try:
        return f"{float(v):.{digits}f}"
    except Exception:
        return "N/A"


def _ensure_dir(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)


def _add_picture_slide(prs: Presentation, title: str, image_path: Path, caption: Optional[str] = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    # Add image
    left = Inches(0.6)
    top = Inches(1.4)
    width = Inches(12.0)
    slide.shapes.add_picture(str(image_path), left, top, width=width)
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.6))
        tx.text_frame.text = caption
        for p in tx.text_frame.paragraphs:
            p.font.size = Pt(12)


def _add_metrics_table_slide(prs: Presentation, title: str, eval_results: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows = 4
    cols = 4
    left = Inches(0.6)
    top = Inches(1.5)
    width = Inches(12.0)
    height = Inches(2.2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["Model", "F1-micro", "F1-macro", "AUC-macro"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h

    def fill_row(i: int, name: str, obj: Dict[str, Any]) -> None:
        table.cell(i, 0).text = name
        table.cell(i, 1).text = _fmt_float(obj.get("f1_micro"))
        table.cell(i, 2).text = _fmt_float(obj.get("f1_macro"))
        table.cell(i, 3).text = _fmt_float(obj.get("auc_macro"))

    fill_row(1, "TF-IDF+LR", eval_results.get("baseline_tfidf_lr", {}))
    fill_row(2, "Detoxify", eval_results.get("baseline_detoxify_unbiased", {}))
    fill_row(3, "Fine-tuned", eval_results.get("finetuned_transformer", {}))

    for r in range(rows):
        for c in range(cols):
            for p in table.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(14 if r == 0 else 12)


def generate_slides_pptx(
    *,
    bundle: ArtifactBundle,
    out_path: Path,
    title: str = "Toxicity Moderation System",
    subtitle: str = "Auto-generated slides from artifacts",
    author_line: str = "(fill) Student / Team",
) -> Path:
    """Generate a 10–15 slide PPTX deck.

    Robustness: missing artifacts create placeholder slides, not failures.
    """
    _ensure_dir(out_path)

    prs = Presentation()
    tmp_dir = out_path.parent / "_figures"

    _add_title_slide(prs, title, f"{subtitle}\n{author_line}")

    # 2) Business problem
    _add_bullets_slide(
        prs,
        "Business Problem & Motivation",
        [
            "Manual moderation is costly and slow.",
            "Toxicity/hate speech harms user experience and brand trust.",
            "Goal: automate flagging/blocking and route uncertain cases to human review.",
        ],
    )

    # 3) Success metrics
    eval_res = bundle.eval_results or {}
    bench = bundle.benchmark or {}
    fin = eval_res.get("finetuned_transformer", {}) if isinstance(eval_res, dict) else {}
    stats = bench.get("stats", {}) if isinstance(bench, dict) else {}
    _add_bullets_slide(
        prs,
        "Success Metrics (Tech + Business)",
        [
            f"F1-micro (fine-tuned): {_fmt_float(fin.get('f1_micro'))}",
            f"AUC-macro (fine-tuned): {_fmt_float(fin.get('auc_macro'))}",
            f"Latency p95 (agent): {_fmt_float(stats.get('p95_ms'), 2)} ms",
            "Business KPIs (fill): time saved, reduced harmful exposure, fewer escalations.",
        ],
    )

    # 4) Data overview
    train_cfg = bundle.train_config or {}
    ds_cfg = train_cfg.get("dataset", {}) if isinstance(train_cfg, dict) else {}
    _add_bullets_slide(
        prs,
        "Data Overview",
        [
            f"Dataset: {_safe(ds_cfg.get('name'))}",
            f"Labels: {_safe(ds_cfg.get('label_fields'))}",
            "Preprocessing: normalization + optional PII redaction.",
            "Split: train/val/test (deterministic seed).",
        ],
    )

    # 5) Architecture diagram
    try:
        arch_img = tmp_dir / "architecture.png"
        save_architecture_diagram(arch_img)
        _add_picture_slide(prs, "System Architecture", arch_img, caption="Agent routes messages to models + policy decisions + logging.")
    except Exception as e:
        logger.warning(f"Architecture diagram failed: {e}")
        _add_bullets_slide(prs, "System Architecture", ["[diagram unavailable]", "See report for details."])

    # 6) Models
    _add_bullets_slide(
        prs,
        "Models",
        [
            "Baseline 1: TF-IDF + Logistic Regression (fast, interpretable)",
            "Baseline 2: Detoxify (off-the-shelf toxicity model)",
            "Final: Fine-tuned transformer (DeBERTa/RoBERTa) for best accuracy",
        ],
    )

    # 7) Training & tuning
    meta = bundle.model_metadata or {}
    tune = bundle.tune_results or {}
    best = (tune.get("best", {}) if isinstance(tune, dict) else {}) or {}
    params = best.get("params") if isinstance(best, dict) else None
    _add_bullets_slide(
        prs,
        "Training & Optimization",
        [
            f"Backbone: {_safe(meta.get('model_name'))}",
            f"Run ID: {_safe(meta.get('run_id'))}",
            f"Best tuned params: {_safe(json.dumps(params) if params else None)}",
            "Techniques: bf16 on H100, gradient checkpointing, early stopping, pos_weight.",
        ],
    )

    # 8) Evaluation table
    if eval_res:
        _add_metrics_table_slide(prs, "Evaluation Results", eval_res)
    else:
        _add_bullets_slide(prs, "Evaluation Results", ["No eval_results.json found. Run evaluation."])

    # 9) Eval charts
    try:
        eval_img = tmp_dir / "eval_comparison.png"
        p = save_eval_comparison_chart(eval_res, eval_img) if eval_res else None
        if p:
            _add_picture_slide(prs, "F1 Comparison", p)
    except Exception as e:
        logger.warning(f"Eval chart failed: {e}")

    try:
        auc_img = tmp_dir / "auc_per_label.png"
        p = save_auc_per_label_chart(eval_res, auc_img) if eval_res else None
        if p:
            _add_picture_slide(prs, "ROC-AUC per label", p)
    except Exception as e:
        logger.warning(f"AUC per label chart failed: {e}")

    # 10) Agentic component
    try:
        flow_img = tmp_dir / "agent_flow.png"
        save_agent_flow_diagram(flow_img)
        _add_picture_slide(prs, "Agentic Moderation Flow", flow_img)
    except Exception:
        _add_bullets_slide(
            prs,
            "Agentic Moderation Flow",
            [
                "1) Detect language",
                "2) Fast model scoring",
                "3) Second-pass model if needed",
                "4) Apply policy thresholds",
                "5) Human review escalation",
            ],
        )

    # 11) Deployment & latency
    if bench:
        try:
            lat_img = tmp_dir / "latency.png"
            p = save_latency_chart(bench, lat_img)
            if p:
                _add_picture_slide(prs, "Deployment Latency", p, caption="End-to-end agent latency on synthetic safe texts.")
            else:
                raise RuntimeError("No chart")
        except Exception:
            _add_bullets_slide(
                prs,
                "Deployment & Latency",
                [
                    f"p50={_fmt_float(stats.get('p50_ms'),2)} ms",
                    f"p95={_fmt_float(stats.get('p95_ms'),2)} ms",
                    f"p99={_fmt_float(stats.get('p99_ms'),2)} ms",
                    "Deploy as FastAPI REST service (Docker optional).",
                ],
            )
    else:
        _add_bullets_slide(prs, "Deployment & Latency", ["No benchmark found.", "Deploy as FastAPI REST service."])

    # 12) Ethics / privacy / fairness
    fair = bundle.fairness or {}
    if fair:
        try:
            fair_img = tmp_dir / "fairness_gap_toxic.png"
            p = save_fairness_gap_chart(fair, fair_img, label="toxic")
            if p:
                _add_picture_slide(prs, "Fairness Slice Evaluation", p, caption="Gaps vs no-identity slice (heuristic identity terms).")
            else:
                raise RuntimeError("No fairness chart")
        except Exception as e:
            logger.warning(f"Fairness chart failed: {e}")
            _add_bullets_slide(
                prs,
                "Ethics, Privacy & Fairness",
                [
                    "Bias risk: identity mentions may be over-flagged.",
                    "Mitigation: monitor slice gaps, adjust thresholds, add in-domain data.",
                    "Privacy: avoid storing raw toxic text; redact PII in logs.",
                ],
            )
    else:
        _add_bullets_slide(
            prs,
            "Ethics, Privacy & Fairness",
            [
                "Bias risk: identity mentions may be over-flagged.",
                "Mitigation: monitoring + human review + domain data.",
                "Privacy: avoid storing raw toxic text; redact PII in logs.",
            ],
        )

    # 13) Key takeaways
    _add_bullets_slide(
        prs,
        "Key Takeaways & Next Steps",
        [
            "End-to-end, deployable moderation system with agentic routing.",
            "Baselines + fine-tuning + monitoring, error analysis, fairness evaluation.",
            "Next: active learning loop, drift detection dashboard, more languages/domains.",
        ],
    )

    prs.save(str(out_path))
    logger.info(f"Generated slide deck at {out_path}")
    return out_path
